"""
PTD-BR Corpus — Gerador de Apresentação PowerPoint
IPEA / COGIT / DIEST
Uso: python gerar_apresentacao.py
Output: apresentacao_ptd_corpus_v4.pptx

v4.0 (2026-04-06): 13 slides, KPIs de 18.4K linhas/59 órgãos,
  +slide Evolução (155 iterações), +slide Balanço S5,
  +slide Arquitetura de Aprendizado VSM
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ── Paleta IPEA ───────────────────────────────────────
AZUL_IPEA   = RGBColor(0x0D, 0x2B, 0x4E)   # azul escuro
VERDE_IPEA  = RGBColor(0x1A, 0x7F, 0x7A)   # verde-água
LARANJA     = RGBColor(0xD9, 0x77, 0x06)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF4, 0xF4, 0xF4)
CINZA_TEXTO = RGBColor(0x44, 0x44, 0x44)

def nova_apresentacao():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def fundo_colorido(slide, cor: RGBColor):
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = cor
    bg.line.fill.background()
    bg.zorder = 0

def caixa(slide, texto, l, t, w, h, tam=18, bold=False,
          cor_txt=BRANCO, cor_bg=None, alinha=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alinha
    run = p.add_run()
    run.text = texto
    run.font.size  = Pt(tam)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = cor_txt
    if cor_bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = cor_bg
    return txb

def retangulo(slide, l, t, w, h, cor: RGBColor, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = cor
    s.line.fill.background()
    return s

def linha_div(slide, t, cor=VERDE_IPEA):
    retangulo(slide, 0.5, t, 12.3, 0.04, cor)

# ════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ════════════════════════════════════════════════════════
def slide_capa(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, AZUL_IPEA)
    retangulo(sl, 0, 5.8, 13.33, 1.7, VERDE_IPEA)
    caixa(sl, 'PTD-BR CORPUS', 0.8, 0.6, 11, 1.2, tam=40, bold=True)
    caixa(sl, 'Planos de Transformação Digital do Governo Federal Brasileiro',
          0.8, 1.7, 11, 0.8, tam=20)
    caixa(sl, 'Coleta · Extração · Estruturação · Análise',
          0.8, 2.35, 10, 0.6, tam=15, italic=True)
    linha_div(sl, 3.1, LARANJA)
    caixa(sl, 'IPEA / COGIT / DIEST', 0.8, 3.3, 8, 0.5, tam=13)
    caixa(sl, 'Denise do Carmo Direito  ·  Lucas Freire Silva',
          0.8, 3.75, 10, 0.5, tam=13)
    caixa(sl, f'Versão 3.0  ·  {datetime.date.today().strftime("%d/%m/%Y")}',
          0.8, 5.95, 10, 0.5, tam=13, bold=True)
    caixa(sl, 'Documento de trabalho — não citar sem autorização dos autores',
          0.8, 6.5, 11.5, 0.4, tam=10, italic=True,
          cor_txt=RGBColor(0xCC,0xCC,0xCC))

# ════════════════════════════════════════════════════════
# SLIDE 2 — O que são os PTDs
# ════════════════════════════════════════════════════════
def slide_contexto(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'O QUE SÃO OS PTDs?', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    itens = [
        ('Base legal', 'Decreto nº 12.198/2024 · Portaria SGD/MGI nº 6.618/2024'),
        ('Obrigação', '90 órgãos da APF pactuam metas digitais com o MGI para 2025–2027'),
        ('Estrutura', '6 Eixos da EFGD (Estratégia Federal de Governo Digital)'),
        ('Documentos', 'Anexo de Entregas (metas) + Documento Diretivo (riscos e diretrizes)'),
        ('Publicação', 'Portal Gov Digital — PDFs públicos assinados digitalmente'),
    ]
    for i, (titulo, desc) in enumerate(itens):
        y = 1.3 + i * 1.1
        retangulo(sl, 0.5, y, 0.08, 0.55, VERDE_IPEA)
        caixa(sl, titulo, 0.75, y, 2.8, 0.35, tam=13, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, desc,   0.75, y+0.32, 11.5, 0.5, tam=12, cor_txt=CINZA_TEXTO)

    caixa(sl, 'Por que coletar?', 0.5, 6.6, 5, 0.4, tam=11, bold=True,
          italic=True, cor_txt=AZUL_IPEA)
    caixa(sl, 'Nenhuma base estruturada dos PTDs existia — este corpus é inédito.',
          3.2, 6.6, 9.5, 0.4, tam=11, italic=True, cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 3 — Arquitetura do Pipeline
# ════════════════════════════════════════════════════════
def slide_arquitetura(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'ARQUITETURA DO PIPELINE', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    camadas = [
        (AZUL_IPEA,  'LAYER 1', 'Coleta',
         'Scraping dinâmico do portal → URLs dos PDFs → Download com SHA-256'),
        (VERDE_IPEA, 'LAYER 2', 'Extração',
         'Docling TableFormer → tabelas dos PDFs → corpus_raw.csv + manifesto'),
        (LARANJA,    'LAYER 3', 'Curadoria',
         'Parser texto→campos · Correção de eixo · Tipo de entrega · Flag IA'),
        (RGBColor(0x45,0x7B,0x9D), 'LAYER 4', 'Análise',
         'Visualizações · Relatório HTML · Metadados JSON · Apresentação'),
    ]
    for i, (cor, label, nome, desc) in enumerate(camadas):
        y = 1.25 + i * 1.45
        retangulo(sl, 0.5, y, 1.6, 1.1, cor)
        caixa(sl, label, 0.5, y + 0.1, 1.6, 0.4, tam=10, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, nome,  0.5, y + 0.5, 1.6, 0.45, tam=16, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, desc,  2.4, y + 0.25, 10.4, 0.6, tam=13, cor_txt=CINZA_TEXTO)
        if i < 3:
            caixa(sl, '▼', 0.9, y + 1.15, 1, 0.25, tam=14, cor_txt=cor)

# ════════════════════════════════════════════════════════
# SLIDE 4 — O que foi coletado
# ════════════════════════════════════════════════════════
def slide_coleta(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'O QUE FOI COLETADO', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    n_reg  = f'{_KPIS["n_registros"]:,}'.replace(',', '.')
    n_orgs = str(_KPIS['n_orgaos'])
    n_serv = f'{_KPIS["n_servicos"]:,}'.replace(',', '.')
    n_risk = str(_KPIS['n_riscos'])
    numeros = [
        (n_orgs, 'órgãos\ncobertos'),
        ('~130', 'PDFs\nbaixados'),
        (n_reg,  'comprometimentos\nno corpus'),
        (n_risk, 'riscos\nextraídos'),
    ]
    for i, (num, leg) in enumerate(numeros):
        x = 0.7 + i * 3.1
        retangulo(sl, x, 1.35, 2.6, 1.7, AZUL_IPEA)
        caixa(sl, num, x, 1.4, 2.6, 1.0, tam=36 if len(num)>5 else 40, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, leg, x, 2.4, 2.6, 0.6, tam=11,
              alinha=PP_ALIGN.CENTER, italic=True)

    linha_div(sl, 3.4)
    caixa(sl, f'Serviços únicos: {n_serv}  ·  Fator de multiplicidade: {_KPIS["fator_mult"]:.2f} prod/serv',
          0.5, 3.55, 12.3, 0.4, tam=13, bold=True, cor_txt=AZUL_IPEA)
    caixa(sl, 'Grupos compartilhados (PDF único, N órgãos):',
          0.5, 4.05, 12, 0.35, tam=12, bold=True, cor_txt=AZUL_IPEA)
    grupos = [
        'MMA → IBAMA · ICMBio · SFB · JBRJ  (1.715 comprometimentos)',
        'MF  → RFB · PGFN · STN  (993 comprometimentos)',
        'MT  → ANTT · DNIT  (132 comprometimentos)',
        'MDA → CONAB  (96 comprometimentos)',
    ]
    for i, g in enumerate(grupos):
        caixa(sl, f'• {g}', 0.9, 4.45 + i*0.48, 12, 0.4, tam=11, cor_txt=CINZA_TEXTO)

    caixa(sl, 'Extração autônoma: 155 iterações do loop watcher · Docling TableFormer · pytesseract fallback',
          0.5, 6.55, 12.3, 0.5, tam=10, italic=True, cor_txt=LARANJA)

# ════════════════════════════════════════════════════════
# SLIDE 5 — Corpus de Entregas
# ════════════════════════════════════════════════════════
def slide_corpus(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    caixa(sl, 'CORPUS DE ENTREGAS', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    ep = _KPIS['eixo_pcts']
    n_reg = f'{_KPIS["n_registros"]:,}'.replace(',','.')
    caixa(sl, 'Unidade: comprometimento = serviço × produto  ·  '
              f'{n_reg} linhas  ·  {_KPIS["n_servicos"]:,} serviços únicos'.replace(',','.'),
          0.5, 1.2, 12.3, 0.5, tam=13, bold=True, cor_txt=AZUL_IPEA,
          alinha=PP_ALIGN.CENTER)
    caixa(sl, f'Fator de multiplicidade: {_KPIS["fator_mult"]:.2f} produtos/serviço  ·  '
              f'Template SGD v2.3 · Loop autônomo ativo (iter. {_KPIS["iteration"]})',
          0.5, 1.65, 12.3, 0.5, tam=11, italic=True, cor_txt=CINZA_TEXTO,
          alinha=PP_ALIGN.CENTER)

    def _ep(k, default):
        v = ep.get(str(k), default)
        return f'{float(v):.1f}%'
    eixos = [
        ('E1', 'Centrado no Cidadão',      _ep(1,'67.7'), AZUL_IPEA),
        ('E2', 'Integrado e Colaborativo', _ep(2,'1.9'),  VERDE_IPEA),
        ('E3', 'Inteligente e Inovador',   _ep(3,'12.4'), LARANJA),
        ('E4', 'Confiável e Seguro',       _ep(4,'17.3'), RGBColor(0xE6,0x39,0x46)),
        ('E5', 'Transparente e Aberto',    _ep(5,'0.2'),  RGBColor(0x45,0x7B,0x9D)),
        ('E6', 'Eficiente e Sustentável',  _ep(6,'0.6'),  RGBColor(0x6A,0x4C,0x93)),
    ]
    for i, (cod, nome, pct, cor) in enumerate(eixos):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2.45 + row * 1.6
        retangulo(sl, x, y, 0.55, 1.1, cor)
        caixa(sl, cod, x, y+0.1, 0.55, 0.5, tam=14, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, pct, x, y+0.6, 0.55, 0.4, tam=11,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, nome, x+0.7, y+0.3, 3.4, 0.6, tam=12, cor_txt=CINZA_TEXTO)

    caixa(sl, '⚠ E1 dominante (67.7%) — possível bias de detecção. E5/E6 sub-representados (<1%).',
          0.5, 6.9, 12.3, 0.35, tam=10, italic=True, cor_txt=LARANJA)

# ════════════════════════════════════════════════════════
# SLIDE 6 — Problemas identificados e corrigidos
# ════════════════════════════════════════════════════════
def slide_problemas(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, RGBColor(0x8B,0x1A,0x1A))
    caixa(sl, 'PROBLEMAS IDENTIFICADOS E CORREÇÕES', 0.5, 0.2, 12, 0.7,
          tam=22, bold=True, alinha=PP_ALIGN.CENTER)

    problemas = [
        ('🔴 Crítico', 'State machine de eixo sem reset por página',
         'Contaminou 21 registros PGFN com E3 incorreto',
         '✅ FIX v3.0: eixo_atual = None a cada nova página'),
        ('🔴 Crítico', 'URLs hardcoded no CATÁLOGO (100/120 com 404)',
         'Pipeline de extração falhava na maioria dos PDFs',
         '✅ FIX v3.0: scraping dinâmico do portal HTML'),
        ('🟡 Médio', 'PRODUTOS hardcoded no script',
         'Exige edição do código a cada novo ciclo PTD',
         '✅ FIX v2.1: externalizado em config/produtos_sgd_v23.json'),
        ('🟡 Médio', 'Correção E3 PGFN-específica (frágil)',
         'Falha para qualquer novo órgão com padrão similar',
         '✅ FIX v2.1: regras em config/correcoes_eixo.json'),
        ('🟠 Menor', 'Sem SHA-256 dos PDFs nos registros',
         'Impossível rastrear registro → PDF de origem',
         '✅ FIX v3.0: pdf_sha256 em cada linha + pipeline_manifest.json'),
    ]
    for i, (sev, prob, impacto, fix) in enumerate(problemas):
        y = 1.2 + i * 1.2
        cor = RGBColor(0x8B,0x1A,0x1A) if 'Crítico' in sev else \
              RGBColor(0xD9,0x77,0x06) if 'Médio' in sev else \
              RGBColor(0xCC,0x88,0x00)
        retangulo(sl, 0.5, y, 0.08, 0.9, cor)
        caixa(sl, prob,    0.75, y,      8.5, 0.38, tam=12, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, impacto, 0.75, y+0.35, 8.5, 0.3,  tam=10, italic=True, cor_txt=CINZA_TEXTO)
        caixa(sl, fix,     9.5,  y+0.15, 3.3, 0.5,  tam=10, cor_txt=VERDE_IPEA, bold=True)

# ════════════════════════════════════════════════════════
# SLIDE 7 — Qualidade do Corpus
# ════════════════════════════════════════════════════════
def slide_qualidade(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'QUALIDADE DO CORPUS', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    pct_ok_str = f'{_KPIS["parse_ok_pct"]:.1f}%'
    sem_prod_n = f'{_KPIS["sem_produto_n"]:,}'.replace(',','.')
    metricas = [
        ('Cobertura de órgãos',   f'{_KPIS["n_orgaos"]}/{_KPIS["n_orgaos"]}',
         'Todos os 59 órgãos com ao menos 1 entrega extraída'),
        ('Parse OK (serv+prod)',   pct_ok_str,
         f'{_KPIS["n_registros"]-_KPIS["sem_produto_n"]:,} linhas com campos identificados'.replace(',','.')),
        ('Cobertura serviço',      f'{_KPIS["cob_servico_pct"]:.1f}%',
         'Serviços identificados no corpus'),
        ('Cobertura produto',      f'{_KPIS["cob_produto_pct"]:.1f}%',
         f'Meta: ≥90% | {sem_prod_n} linhas sem_produto pendentes'),
        ('IA real identificada',   str(_KPIS['ia_real_n']),
         f'Registros em {_KPIS["ia_real_orgs"]} órgãos (PAT_IA_REAL preciso)'),
        ('Cobertura data PTD',     f'{_KPIS["cob_data_pct"]:.1f}%',
         'Datas de entrega extraídas'),
    ]
    for i, (label, valor, desc) in enumerate(metricas):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = 1.3 + row * 1.9
        retangulo(sl, x, y, 6.0, 1.6, BRANCO)
        retangulo(sl, x, y, 0.08, 1.6, VERDE_IPEA)
        caixa(sl, valor, x+0.25, y+0.1, 5.5, 0.8, tam=30, bold=True,
              cor_txt=AZUL_IPEA)
        caixa(sl, label, x+0.25, y+0.85, 5.5, 0.4, tam=12, bold=True,
              cor_txt=CINZA_TEXTO)
        caixa(sl, desc,  x+0.25, y+1.2,  5.5, 0.3, tam=10, italic=True,
              cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 8 — Avaliação de Prontidão
# ════════════════════════════════════════════════════════
def slide_prontidao(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    caixa(sl, 'AVALIAÇÃO DE PRONTIDÃO', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    dimensoes = [
        ('Completude da coleta',     3, 5, 'Scraping dinâmico implementado; ~30% PDFs ainda sem extração validada'),
        ('Qualidade da extração',    4, 5, 'Docling TableFormer + fix state machine; OCR para PDFs tarjados'),
        ('Curadoria semântica',      4, 5, 'Parser, tipo_entrega, ia_real, correção E3 com regras externas'),
        ('Rastreabilidade',          4, 5, 'SHA-256 PDFs + hashes CSV inputs + pipeline_manifest.json'),
        ('Reprodutibilidade',        3, 5, 'Checkpoints implementados; falta documentação de ambiente (requirements.txt)'),
        ('Ciência aberta',           3, 5, 'Código versionado no GitHub; falta DOI de dados e licença explícita'),
        ('Pronto para publicação',   2, 5, 'Requires: validação manual amostra + relatório HTML + DOI Zenodo'),
    ]

    for i, (dim, nota, max_n, obs) in enumerate(dimensoes):
        y = 1.3 + i * 0.84
        # barra de fundo
        retangulo(sl, 3.8, y+0.15, 5.5, 0.45, RGBColor(0xDD,0xDD,0xDD))
        # barra de nota
        cor_bar = VERDE_IPEA if nota >= 4 else LARANJA if nota == 3 else \
                  RGBColor(0xE6,0x39,0x46)
        retangulo(sl, 3.8, y+0.15, 5.5*(nota/max_n), 0.45, cor_bar)
        caixa(sl, dim,    0.5, y, 3.2, 0.5, tam=12, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, f'{nota}/{max_n}', 9.5, y, 0.7, 0.5, tam=14, bold=True,
              cor_txt=cor_bar)
        caixa(sl, obs, 10.3, y+0.05, 2.8, 0.5, tam=9, italic=True,
              cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 9 — Próximos Passos
# ════════════════════════════════════════════════════════
def slide_proximos_passos(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'PRÓXIMOS PASSOS', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    itr = _KPIS['iteration']
    passos = [
        (f'Curto prazo\n(iters {itr}–{itr+19})', VERDE_IPEA, [
            f'pct_ok: {_KPIS["parse_ok_pct"]:.1f}% → 80% via vocab por sigla',
            'INCRA col_map_ok: 0% → >50% (col_keys_extra.json)',
            'Meta-insight #1 gerado na iter. 164',
        ]),
        (f'Médio prazo\n(iters {itr+20}–{itr+39})', LARANJA, [
            'pct_ok ≥ 80% → Stage 2 gate (aprovação humana)',
            'Eixos 2/5/6 > 3% cada (regex expandida)',
            'Publicar corpus no Repositório de Dados IPEA (DOI)',
        ]),
        (f'Longo prazo\n(iters {itr+40}–204)', RGBColor(0x45,0x7B,0x9D), [
            'pct_ok ≥ 88% · learning_signals com 50 entradas',
            'Análise comparativa PTDs 2025 vs 2026',
            'S4 completo: ptd_learning_signals → recomendações automáticas',
        ]),
    ]

    for col, (periodo, cor, items) in enumerate(passos):
        x = 0.5 + col * 4.2
        retangulo(sl, x, 1.3, 3.9, 0.65, cor)
        caixa(sl, periodo, x, 1.35, 3.9, 0.55, tam=13, bold=True,
              alinha=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            caixa(sl, f'• {item}', x+0.1, 2.1 + j*1.5, 3.7, 1.3,
                  tam=11, cor_txt=CINZA_TEXTO)

# ════════════════════════════════════════════════════════
# SLIDE 10 — Balanço de Sessão
# ════════════════════════════════════════════════════════
def slide_balanco(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, AZUL_IPEA)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    caixa(sl, 'BALANÇO DE SESSÃO', 0.5, 0.2, 12, 0.7,
          tam=24, bold=True, alinha=PP_ALIGN.CENTER)

    itr = _KPIS['iteration']
    n_reg = f'{_KPIS["n_registros"]:,}'.replace(',','.')
    entregues = [
        f'Corpus PTD-BR v2.1     — {n_reg} comprometimentos · 59 órgãos · 420 riscos',
        'ptd_pipeline_v30.py     — Pipeline autônomo com loop watcher (155 iterações)',
        'meta_learning.py        — S4 mínimo viável: slope analysis + strategy switch',
        'config/col_keys_extra.json   — Col_map por sigla (Step B watcher)',
        'ptd_constants.py        — Regex eixos 2/5/6 expandida (3× cobertura)',
        'gerar_apresentacao.py   — Este PPT (reproduzível) + gerar_relatorio_tecnico.py',
    ]
    caixa(sl, '✅ CORPUS E INFRAESTRUTURA', 0.5, 1.2, 12, 0.4,
          tam=14, bold=True, cor_txt=VERDE_IPEA)
    for i, e in enumerate(entregues):
        caixa(sl, e, 0.7, 1.65 + i*0.45, 12, 0.4, tam=11)

    linha_div(sl, 4.5, LARANJA)

    caixa(sl, '⚠️  PENDENTE (próximas 50 iterações — iter. 155→204)', 0.5, 4.65, 12, 0.4,
          tam=14, bold=True, cor_txt=LARANJA)
    pendentes = [
        f'pct_ok: {_KPIS["parse_ok_pct"]:.1f}% → meta 90% (via vocab + col_keys expansão autônoma)',
        'INCRA col_map_ok: 0% → >50% (headers capturados, col_keys_extra.json no próximo run)',
        'Eixos 2/5/6: <1% → esperado >3% (regex expandida no último commit)',
    ]
    for i, p in enumerate(pendentes):
        caixa(sl, f'• {p}', 0.7, 5.1 + i*0.43, 12, 0.4,
              tam=11, cor_txt=LARANJA)

    caixa(sl, f'Gerado em {datetime.datetime.now().strftime("%d/%m/%Y %H:%M")} | COGIT/DIEST/IPEA',
          0.5, 7.1, 12.3, 0.3, tam=9, italic=True,
          cor_txt=RGBColor(0xAA,0xAA,0xAA))

# ════════════════════════════════════════════════════════
# SLIDE 5b — Evolução: 155 iterações
# ════════════════════════════════════════════════════════
def slide_evolucao(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, VERDE_IPEA)
    itr = _KPIS['iteration']
    caixa(sl, f'EVOLUÇÃO — {itr} ITERAÇÕES AUTÔNOMAS', 0.5, 0.2, 12, 0.7,
          tam=22, bold=True, alinha=PP_ALIGN.CENTER)

    # Timeline horizontal
    marcos = [
        (0,    'Stage 0\nCobertura', '10 órgãos\nzerados',   RGBColor(0xE6,0x39,0x46)),
        (50,   'Stage 0 →1\nConcluído', '59/59\nórgãos OK',  VERDE_IPEA),
        (100,  'Stage 1\nQualidade', 'vocab loop\nvazio',     LARANJA),
        (154,  'S4 ativo\nmeta_learning', 'unmatched\npor sigla', AZUL_IPEA),
        (itr,  f'Iter. {itr}\nAtual', 'pct_ok=68%\nstage=1',  VERDE_IPEA),
    ]
    # Linha base
    retangulo(sl, 0.8, 3.5, 11.5, 0.06, RGBColor(0xBB,0xBB,0xBB))
    max_iter = max(itr, 160)
    for it, titulo, desc, cor in marcos:
        x = 0.8 + (it / max_iter) * 11.3
        retangulo(sl, x-0.08, 3.1, 0.16, 0.7, cor)
        caixa(sl, titulo, x-1.0, 1.9, 2.0, 0.9, tam=10, bold=True,
              cor_txt=cor, alinha=PP_ALIGN.CENTER)
        caixa(sl, desc, x-1.0, 4.0, 2.0, 0.8, tam=9, italic=True,
              cor_txt=CINZA_TEXTO, alinha=PP_ALIGN.CENTER)

    linha_div(sl, 5.1)
    # Métricas de evolução
    cols = [
        ('Stage 0', 'Concluído\niter. ~50', VERDE_IPEA),
        ('Stage 1', f'Ativo\niter. {itr}', LARANJA),
        ('pct_ok', f'{_KPIS["parse_ok_pct"]:.1f}%\n(meta: 90%)', AZUL_IPEA),
        ('S4 L-C', 'learning_signals\nativo', VERDE_IPEA),
        ('Meta-learn', 'a cada\n10 iterações', LARANJA),
    ]
    for i, (tit, val, cor) in enumerate(cols):
        x = 0.7 + i * 2.4
        retangulo(sl, x, 5.35, 2.1, 1.5, cor)
        caixa(sl, tit, x, 5.4, 2.1, 0.4, tam=11, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, val, x, 5.8, 2.1, 0.9, tam=12,
              alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 6b — Balanço S5
# ════════════════════════════════════════════════════════
def slide_balanco_s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'BALANÇO S5 — ATENDIMENTO DO PROPÓSITO', 0.5, 0.2, 12, 0.7,
          tam=22, bold=True, alinha=PP_ALIGN.CENTER)

    # 3 critérios S5
    criterios = [
        ('pct_ok ≥ 90%',
         f'{_KPIS["parse_ok_pct"]:.1f}% atual',
         f'−{90-_KPIS["parse_ok_pct"]:.1f}pp',
         '❌',
         'Vocab insuficiente (54 produtos), col_map=0% INCRA, watcher loop vazio (resolvido iter. 155)'),
        ('Riscos ≥ 80% órgãos',
         f'31/59 órgãos (52%)',
         '−16 órgãos',
         '❌',
         'Extrator L7 funciona para 31 órgãos; 28 sem riscos (ausência real ou falha do parser)'),
        ('Eixos 2/5/6 > 5% cada',
         'E2=1.9%  E5=0.2%  E6=0.6%',
         '< 1% em 3 eixos',
         '❌',
         'Regex expandida (iter. 155) — resultado verificável no próximo run completo'),
    ]
    for i, (criterio, atual, delta, status, causa) in enumerate(criterios):
        y = 1.25 + i * 1.85
        cor_st = RGBColor(0xE6,0x39,0x46) if status == '❌' else VERDE_IPEA
        retangulo(sl, 0.5, y, 12.33, 1.6, BRANCO)
        retangulo(sl, 0.5, y, 0.08, 1.6, cor_st)
        caixa(sl, criterio, 0.75, y+0.05, 3.8, 0.5, tam=13, bold=True, cor_txt=AZUL_IPEA)
        caixa(sl, atual,    0.75, y+0.5,  3.8, 0.4, tam=12, cor_txt=CINZA_TEXTO)
        caixa(sl, delta,    4.8,  y+0.2,  1.8, 0.5, tam=14, bold=True, cor_txt=cor_st,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, status,   6.8,  y+0.2,  0.8, 0.6, tam=24, bold=True, cor_txt=cor_st,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, f'Causa: {causa}', 7.7, y+0.1, 4.8, 1.4, tam=9,
              italic=True, cor_txt=CINZA_TEXTO)

    linha_div(sl, 6.9)
    caixa(sl, f'Critérios atendidos: 0 de 3  ·  Meta de convergência: iter. 204 (pct_ok ≥ 88%)',
          0.5, 7.0, 12.3, 0.35, tam=11, bold=True, cor_txt=LARANJA,
          alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE 9b — Arquitetura de Aprendizado VSM
# ════════════════════════════════════════════════════════
def slide_vsm_learning(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_colorido(sl, CINZA_CLARO)
    retangulo(sl, 0, 0, 13.33, 1.1, AZUL_IPEA)
    caixa(sl, 'ARQUITETURA DE APRENDIZADO — VSM CIBERNÉTICO', 0.5, 0.2, 12, 0.7,
          tam=22, bold=True, alinha=PP_ALIGN.CENTER)

    # S1-S5 coluna esquerda
    sistemas = [
        ('S5', 'Política',      'Operador humano + thresholds (90% pct_ok, 80% riscos)',    AZUL_IPEA),
        ('S4', 'Inteligência',  'meta_learning.py · ptd_learning_signals.json (S4 ativo!)', VERDE_IPEA),
        ('S3', 'Gestão',        'watcher.yml · state machine · problem classification',     VERDE_IPEA),
        ('S3*','Auditoria',     'gerar_relatorio.py · ptd_run_summary.json · sensor',       LARANJA),
        ('S2', 'Coordenação',   'config/ (vocab, col_keys, normas) — estático hoje',        LARANJA),
        ('S1', 'Operações',     'ptd_pipeline_v30.py · ptd_corpus_v21.py por órgão/PDF',   AZUL_IPEA),
    ]
    for i, (cod, nome, desc, cor) in enumerate(sistemas):
        y = 1.2 + i * 0.97
        retangulo(sl, 0.4, y, 0.55, 0.78, cor)
        caixa(sl, cod, 0.4, y+0.1, 0.55, 0.3, tam=11, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, nome, 1.1, y+0.05, 2.0, 0.35, tam=12, bold=True, cor_txt=cor)
        caixa(sl, desc, 1.1, y+0.38, 5.8, 0.38, tam=9, italic=True, cor_txt=CINZA_TEXTO)

    # Camadas L-A/B/C coluna direita
    linha_div_v = lambda x, t1, t2: retangulo(sl, x, t1, 0.04, t2-t1, RGBColor(0xBB,0xBB,0xBB))
    retangulo(sl, 7.5, 1.2, 5.5, 0.06, RGBColor(0xBB,0xBB,0xBB))  # linha topo
    camadas_l = [
        ('L-C\nglobal',  'ptd_learning_signals.json\nRun history · delta_pct_ok · per_sigla_ok',
         'Acumulado em pipeline-outputs. Trajetória de 50 iters.', AZUL_IPEA),
        ('L-B\nmeta',    'meta_learning.py (iter % 10)\nSlope analysis · strategy switch',
         'slope < 0.1pp/iter → troca vocabulario→col_keys→eixo_regex', VERDE_IPEA),
        ('L-A\niteração','vocab + col_keys + eixo regex\npor-sigla + unrecognized_headers',
         'Watcher Steps A/B/D a cada ciclo de 5 min', LARANJA),
    ]
    for i, (camada, titulo, desc, cor) in enumerate(camadas_l):
        y = 1.35 + i * 1.87
        retangulo(sl, 7.5, y, 0.6, 1.5, cor)
        caixa(sl, camada, 7.5, y+0.3, 0.6, 0.8, tam=9, bold=True,
              alinha=PP_ALIGN.CENTER)
        caixa(sl, titulo, 8.2, y+0.05, 5.0, 0.7, tam=10, bold=True, cor_txt=cor)
        caixa(sl, desc,   8.2, y+0.75, 5.0, 0.65, tam=9, italic=True, cor_txt=CINZA_TEXTO)

    linha_div(sl, 7.1)
    caixa(sl, 'S4 passa de 🔴 para 🟡 ativo: ptd_learning_signals.json acumulando desde iter. 155',
          0.5, 7.15, 12.3, 0.3, tam=10, italic=True, cor_txt=VERDE_IPEA,
          alinha=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# GERAR
# ════════════════════════════════════════════════════════
def _carregar_kpis() -> dict:
    """Lê metadados, pivot e summary para KPIs reais."""
    from pathlib import Path as _P
    import json as _j
    kpis: dict = {}
    meta_path = _P('ptd_corpus/03_database/ptd_corpus_v21_metadados.json')
    if meta_path.exists():
        try:
            m = _j.loads(meta_path.read_text(encoding='utf-8'))
            c = m.get('corpus', {})
            p = m.get('parse', {})
            e = m.get('eixos', {}).get('original', {})
            tot_e = sum(e.values()) or 1
            kpis.update({
                'n_registros':     c.get('total_linhas', 18396),
                'n_orgaos':        m.get('proveniencia', {}).get('n_orgaos', 59),
                'n_servicos':      c.get('servicos_unicos', 3376),
                'fator_mult':      c.get('fator_mult_medio', 5.45),
                'n_riscos':        m.get('proveniencia', {}).get('n_riscos', 420),
                'parse_ok_pct':    round(p.get('ok', 13887) / max(c.get('total_linhas',1),1)*100,1),
                'cob_servico_pct': p.get('cobertura_servico_pct', 97.1),
                'cob_produto_pct': p.get('cobertura_produto_pct', 76.2),
                'cob_data_pct':    p.get('cobertura_data_ptd_pct', 73.9),
                'ia_real_n':       m.get('ia_real', {}).get('registros', 48),
                'ia_real_orgs':    m.get('ia_real', {}).get('orgaos', 11),
                'eixo_pcts': {str(k): round(v/tot_e*100,1) for k,v in e.items()},
                'sem_produto_n':   p.get('sem_produto', 3969),
            })
        except Exception:
            pass
    # Iteration counter from .trigger_debug
    td = _P('.trigger_debug')
    if td.exists():
        for line in td.read_text().splitlines():
            if line.startswith('iteration:'):
                try: kpis['iteration'] = int(line.split(':')[1].strip())
                except: pass
    kpis.setdefault('n_registros',  18396)
    kpis.setdefault('n_orgaos',     59)
    kpis.setdefault('n_servicos',   3376)
    kpis.setdefault('fator_mult',   5.45)
    kpis.setdefault('n_riscos',     420)
    kpis.setdefault('parse_ok_pct', 75.5)
    kpis.setdefault('cob_servico_pct', 97.1)
    kpis.setdefault('cob_produto_pct', 76.2)
    kpis.setdefault('cob_data_pct',    73.9)
    kpis.setdefault('ia_real_n',    48)
    kpis.setdefault('ia_real_orgs', 11)
    kpis.setdefault('iteration',    155)
    kpis.setdefault('eixo_pcts', {'1':'67.7','2':'1.9','3':'12.4','4':'17.3','5':'0.2','6':'0.6'})
    kpis.setdefault('sem_produto_n', 3969)
    return kpis

# KPIs globais — preenchidos com valores reais se disponíveis, senão placeholders
_KPIS = _carregar_kpis()

def gerar():
    from pathlib import Path as _Path
    prs = nova_apresentacao()
    slide_capa(prs)          # 1
    slide_contexto(prs)      # 2
    slide_arquitetura(prs)   # 3
    slide_coleta(prs)        # 4 — atualizado: 59 órgãos, 18.4K linhas
    slide_corpus(prs)        # 5 — atualizado: eixos reais, iter atual
    slide_evolucao(prs)      # 5b NOVO: 155 iterações, timeline
    slide_problemas(prs)     # 6
    slide_balanco_s5(prs)    # 6b NOVO: balanço S5, 3 critérios
    slide_qualidade(prs)     # 7 — atualizado: métricas reais
    slide_prontidao(prs)     # 8
    slide_vsm_learning(prs)  # 9b NOVO: VSM + L-A/B/C
    slide_proximos_passos(prs)  # 9 — atualizado: roadmap 50 iters
    slide_balanco(prs)       # 10 — atualizado: 155 iters, S4 ativo
    out = _Path('apresentacao_ptd_corpus_v4.pptx')
    prs.save(str(out))
    print(f'✅ Apresentação salva: {out}  ({out.stat().st_size//1024} KB)')
    print(f'   {len(prs.slides)} slides')
    return out

if __name__ == '__main__':
    gerar()
